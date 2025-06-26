import os
import json
import uuid
from typing import Dict, List, Any, Optional
import google.generativeai as genai
from sentence_transformers import SentenceTransformer
import chromadb
from chromadb.config import Settings
import fitz  # PyMuPDF
from pptx import Presentation
from PIL import Image
import pytesseract
import base64
import io 
import re
import numpy as np

class EnhancedRAGSystem:
    def __init__(self, vector_db_dir: str, openai_api_key: str = None):
        self.vector_db_dir = vector_db_dir
        
        # Initialize OpenAI
        
        genai.configure(api_key=os.getenv('GOOGLE_API_KEY'))
        self.model = genai.GenerativeModel('gemini-2.0-flash')

        # Initialize ChromaDB
        self.client = chromadb.PersistentClient(
            path=vector_db_dir,
            settings=Settings(anonymized_telemetry=False)
        )
        
        # Initialize embedding model
        self.embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
        
        # Create collections
        self.text_collection = self._get_or_create_collection("text_chunks")
        self.image_collection = self._get_or_create_collection("images")
        self.code_collection = self._get_or_create_collection("code_blocks")
        self.formula_collection = self._get_or_create_collection("formulas")
        
        # Store processed images directory
        self.images_dir = os.path.join(vector_db_dir, "extracted_images")
        os.makedirs(self.images_dir, exist_ok=True)
    
    def _get_or_create_collection(self, name: str):
        """Get or create a ChromaDB collection"""
        try:
            return self.client.get_collection(name)
        except Exception:
            return self.client.create_collection(
                name=name,
                metadata={"hnsw:space": "cosine"}
            )
    
    def process_and_add_document(self, file_path: str, filename: str) -> Dict[str, Any]:
        """Process document with enhanced OCR and content extraction"""
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext == '.pdf':
            processed_content = self._process_pdf_enhanced(file_path, filename)
        elif file_ext in ['.ppt', '.pptx']:
            processed_content = self._process_ppt_enhanced(file_path, filename)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")
        
        # Add to vector database
        doc_id = self._add_to_vector_db(processed_content, filename)
        processed_content["document_id"] = doc_id
        
        return processed_content
    
    def _process_pdf_enhanced(self, file_path: str, filename: str) -> Dict[str, Any]:
        """Enhanced PDF processing with better OCR and content extraction"""
        doc = fitz.open(file_path)
        
        result = {
            "filename": filename,
            "file_type": "pdf",
            "text_chunks": [],
            "images": [],
            "formulas": [],
            "code_blocks": [],
            "metadata": {
                "total_pages": len(doc),
                "file_size": os.path.getsize(file_path)
            }
        }
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # Extract text with better chunking
            text = page.get_text()
            if text.strip():
                chunks = self._create_intelligent_chunks(text, page_num + 1)
                result["text_chunks"].extend(chunks)
                
                # Extract code blocks with better patterns
                code_blocks = self._extract_code_blocks_enhanced(text, page_num + 1)
                result["code_blocks"].extend(code_blocks)
                
                # Extract formulas with better detection
                formulas = self._extract_formulas_enhanced(text, page_num + 1)
                result["formulas"].extend(formulas)
            
            # Enhanced image extraction with OCR
            images = self._extract_images_enhanced(page, page_num + 1, filename)
            result["images"].extend(images)
        
        doc.close()
        return result
    
    def _process_ppt_enhanced(self, file_path: str, filename: str) -> Dict[str, Any]:
        """Enhanced PowerPoint processing"""
        prs = Presentation(file_path)
        
        result = {
            "filename": filename,
            "file_type": "pptx",
            "text_chunks": [],
            "images": [],
            "formulas": [],
            "code_blocks": [],
            "metadata": {
                "total_slides": len(prs.slides),
                "file_size": os.path.getsize(file_path)
            }
        }
        
        for slide_num, slide in enumerate(prs.slides, 1):
            # Extract all text from slide
            slide_text = self._extract_slide_text(slide)
            
            if slide_text.strip():
                chunks = self._create_intelligent_chunks(slide_text, slide_num, is_slide=True)
                result["text_chunks"].extend(chunks)
                
                code_blocks = self._extract_code_blocks_enhanced(slide_text, slide_num, is_slide=True)
                result["code_blocks"].extend(code_blocks)
                
                formulas = self._extract_formulas_enhanced(slide_text, slide_num, is_slide=True)
                result["formulas"].extend(formulas)
            
            # Extract images from slide
            images = self._extract_slide_images_enhanced(slide, slide_num, filename)
            result["images"].extend(images)
        
        return result
    
    def _extract_slide_text(self, slide) -> str:
        """Extract all text from a PowerPoint slide"""
        text_content = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text.strip())
            
            # Handle tables
            if shape.shape_type == 19:  # Table
                try:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            text_content.append(" | ".join(row_text))
                except:
                    pass
        
        return "\n".join(text_content)
    
    def _create_intelligent_chunks(self, text: str, page_num: int, is_slide: bool = False) -> List[Dict[str, Any]]:
        """Create intelligent text chunks with context preservation"""
        chunks = []
        
        # Split by sections first (headers, bullet points, etc.)
        sections = self._identify_text_sections(text)
        
        for i, section in enumerate(sections):
            if len(section.strip()) < 30:  # Skip very short sections
                continue
            
            chunk = {
                "text": section.strip(),
                "chunk_id": f"{'slide' if is_slide else 'page'}_{page_num}_chunk_{i}",
                "page_number" if not is_slide else "slide_number": page_num,
                "chunk_type": "text",
                "word_count": len(section.split()),
                "char_count": len(section),
                "section_type": self._classify_text_section(section)
            }
            chunks.append(chunk)
        
        return chunks
    
    def _identify_text_sections(self, text: str) -> List[str]:
        """Identify logical sections in text"""
        # Split by headers, bullet points, and paragraphs
        import re
        
        # Look for headers (numbered, bulleted, or capitalized)
        header_pattern = r'(?:^|\n)(?:\d+\.|\*|\-|[A-Z][A-Z\s]+:|\#{1,3}\s)'
        sections = re.split(header_pattern, text, flags=re.MULTILINE)
        
        # If no clear sections, split by double newlines
        if len(sections) <= 1:
            sections = [s.strip() for s in text.split('\n\n') if s.strip()]
        
        return [s.strip() for s in sections if s.strip()]
    
    def _classify_text_section(self, text: str) -> str:
        """Classify the type of text section"""
        text_lower = text.lower()
        
        if any(word in text_lower for word in ['definition', 'define', 'what is', 'refers to']):
            return 'definition'
        elif any(word in text_lower for word in ['example', 'for instance', 'such as']):
            return 'example'
        elif any(word in text_lower for word in ['theorem', 'proof', 'lemma', 'proposition']):
            return 'theorem'
        elif any(word in text_lower for word in ['algorithm', 'procedure', 'steps']):
            return 'algorithm'
        else:
            return 'general'
    
    def _extract_code_blocks_enhanced(self, text: str, page_num: int, is_slide: bool = False) -> List[Dict[str, Any]]:
        """Enhanced code block extraction"""
        import re
        
        code_blocks = []
        
        # Enhanced patterns for different programming languages
        patterns = [
            # Python
            (r'def\s+\w+\([^)]*\):\s*\n(?:\s{4}.*\n)*', 'python'),
            (r'class\s+\w+.*?:\s*\n(?:\s{4}.*\n)*', 'python'),
            (r'import\s+[\w.]+|from\s+[\w.]+\s+import', 'python'),
            (r'for\s+\w+\s+in\s+.*?:\s*\n(?:\s{4}.*\n)*', 'python'),
            (r'if\s+.*?:\s*\n(?:\s{4}.*\n)*', 'python'),
            
            # JavaScript
            (r'function\s+\w+\([^)]*\)\s*{[^}]*}', 'javascript'),
            (r'const\s+\w+\s*=\s*.*?=>', 'javascript'),
            (r'let\s+\w+\s*=.*?;', 'javascript'),
            
            # SQL
            (r'SELECT\s+.*?FROM\s+.*?(?:WHERE\s+.*?)?(?:;|$)', 'sql'),
            (r'CREATE\s+TABLE\s+.*?(?:;|$)', 'sql'),
            
            # Mathematical expressions
            (r'\$\$.*?\$\$', 'latex'),
            (r'\$.*?\$', 'latex'),
            
            # Code blocks (markdown style)
            (r'```[\w]*\n.*?\n```', 'code'),
            
            # Indented code blocks
            (r'(?:^|\n)(?: {4}|\t).*(?:\n(?: {4}|\t).*)*', 'code')
        ]
        
        for i, (pattern, language) in enumerate(patterns):
            matches = re.finditer(pattern, text, re.MULTILINE | re.DOTALL | re.IGNORECASE)
            for j, match in enumerate(matches):
                code = match.group().strip()
                if len(code) > 10:  # Filter out very short matches
                    code_block = {
                        "code": code,
                        "code_id": f"{'slide' if is_slide else 'page'}_{page_num}_code_{i}_{j}",
                        "page_number" if not is_slide else "slide_number": page_num,
                        "language": language,
                        "chunk_type": "code",
                        "confidence": self._calculate_code_confidence(code, language)
                    }
                    code_blocks.append(code_block)
        
        return code_blocks
    
    def _calculate_code_confidence(self, code: str, language: str) -> float:
        """Calculate confidence that text is actually code"""
        confidence = 0.5
        
        # Language-specific indicators
        if language == 'python':
            indicators = ['def ', 'class ', 'import ', 'if __name__', 'print(', 'return ']
        elif language == 'javascript':
            indicators = ['function ', 'const ', 'let ', 'var ', 'console.log', '=>']
        elif language == 'sql':
            indicators = ['SELECT', 'FROM', 'WHERE', 'INSERT', 'UPDATE', 'DELETE']
        else:
            indicators = ['(', ')', '{', '}', ';', '=']
        
        for indicator in indicators:
            if indicator in code:
                confidence += 0.1
        
        return min(confidence, 1.0)
    
    def _extract_formulas_enhanced(self, text: str, page_num: int, is_slide: bool = False) -> List[Dict[str, Any]]:
        """Enhanced formula extraction"""
        import re
        
        formulas = []
        
        # Enhanced patterns for mathematical expressions
        patterns = [
            # LaTeX math
            (r'\$\$([^$]+)\$\$', 'latex_display'),
            (r'\$([^$]+)\$', 'latex_inline'),
            
            # Mathematical symbols and expressions
            (r'[a-zA-Z]\s*=\s*[^,\n.]{3,}', 'equation'),
            (r'[∑∏∫∆∇∞≈≠≤≥±×÷√∠∴∵∪∩∈∀∃][^,\n.]{2,}', 'mathematical'),
            
            # Statistical formulas
            (r'(?:mean|std|var|correlation)\s*=\s*[^,\n.]+', 'statistical'),
            
            # ML/AI formulas
            (r'(?:loss|cost|error)\s*=\s*[^,\n.]+', 'ml_formula'),
            (r'y\s*=\s*[^,\n.]+', 'regression'),
            
            # Physics formulas
            (r'[FEPVT]\s*=\s*[^,\n.]+', 'physics')
        ]
        
        for i, (pattern, formula_type) in enumerate(patterns):
            matches = re.finditer(pattern, text, re.IGNORECASE)
            for j, match in enumerate(matches):
                formula_text = match.group().strip()
                if len(formula_text) > 3:
                    formula = {
                        "formula": formula_text,
                        "formula_id": f"{'slide' if is_slide else 'page'}_{page_num}_formula_{i}_{j}",
                        "page_number" if not is_slide else "slide_number": page_num,
                        "chunk_type": "formula",
                        "formula_type": formula_type,
                        "confidence": self._calculate_formula_confidence(formula_text, formula_type)
                    }
                    formulas.append(formula)
        
        return formulas
    
    def _calculate_formula_confidence(self, formula: str, formula_type: str) -> float:
        """Calculate confidence that text is actually a formula"""
        confidence = 0.5
        
        # Check for mathematical symbols
        math_symbols = ['=', '+', '-', '*', '/', '^', '√', '∑', '∏', '∫']
        for symbol in math_symbols:
            if symbol in formula:
                confidence += 0.1
        
        # Check for variables and numbers
        import re
        if re.search(r'[a-zA-Z]', formula) and re.search(r'\d', formula):
            confidence += 0.2
        
        return min(confidence, 1.0)
    
    def _extract_images_enhanced(self, page, page_num: int, filename: str) -> List[Dict[str, Any]]:
        """Enhanced image extraction with better OCR"""
        images = []
        image_list = page.get_images()
        
        for img_index, img in enumerate(image_list):
            try:
                xref = img[0]
                pix = fitz.Pixmap(page.parent, xref)
                
                if pix.n - pix.alpha < 4:  # RGB
                    img_data = pix.tobytes("png")
                    
                    # Generate unique filename
                    img_filename = f"{filename}_page_{page_num}_img_{img_index}.png"
                    img_path = os.path.join(self.images_dir, img_filename)
                    
                    # Save image
                    with open(img_path, "wb") as img_file:
                        img_file.write(img_data)
                    
                    # Enhanced OCR with multiple methods
                    ocr_results = self._enhanced_ocr(img_path)
                    
                    # Analyze image content
                    image_analysis = self._analyze_image_content(img_path, ocr_results['text'])
                    
                    image_info = {
                        "image_path": img_path,
                        "image_id": f"page_{page_num}_img_{img_index}",
                        "page_number": page_num,
                        "ocr_text": ocr_results['text'],
                        "ocr_confidence": ocr_results['confidence'],
                        "chunk_type": "image",
                        "width": pix.width,
                        "height": pix.height,
                        "image_type": image_analysis['type'],
                        "contains_text": len(ocr_results['text'].strip()) > 10,
                        "contains_math": image_analysis['has_math'],
                        "contains_code": image_analysis['has_code'],
                        "description": image_analysis['description']
                    }
                    images.append(image_info)
                
                pix = None
            except Exception as e:
                print(f"Error extracting image {img_index} from page {page_num}: {e}")
        
        return images
    
    def _extract_slide_images_enhanced(self, slide, slide_num: int, filename: str) -> List[Dict[str, Any]]:
        """Enhanced image extraction from PowerPoint slides"""
        images = []
        
        for shape_index, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    img_filename = f"{filename}_slide_{slide_num}_img_{shape_index}.png"
                    img_path = os.path.join(self.images_dir, img_filename)
                    
                    with open(img_path, "wb") as img_file:
                        img_file.write(image.blob)
                    
                    # Enhanced OCR
                    ocr_results = self._enhanced_ocr(img_path)
                    
                    # Analyze image content
                    image_analysis = self._analyze_image_content(img_path, ocr_results['text'])
                    
                    image_info = {
                        "image_path": img_path,
                        "image_id": f"slide_{slide_num}_img_{shape_index}",
                        "slide_number": slide_num,
                        "ocr_text": ocr_results['text'],
                        "ocr_confidence": ocr_results['confidence'],
                        "chunk_type": "image",
                        "image_type": image_analysis['type'],
                        "contains_text": len(ocr_results['text'].strip()) > 10,
                        "contains_math": image_analysis['has_math'],
                        "contains_code": image_analysis['has_code'],
                        "description": image_analysis['description']
                    }
                    images.append(image_info)
                
                except Exception as e:
                    print(f"Error extracting image {shape_index} from slide {slide_num}: {e}")
        
        return images
    
    def _enhanced_ocr(self, image_path: str) -> Dict[str, Any]:
        """Enhanced OCR with multiple methods and confidence scoring"""
        try:
            image = Image.open(image_path)
            
            # Standard OCR
            text = pytesseract.image_to_string(image, config='--psm 6')
            
            # Get confidence data
            data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
            confidences = [int(conf) for conf in data['conf'] if int(conf) > 0]
            avg_confidence = sum(confidences) / len(confidences) if confidences else 0
            
            # Try different PSM modes for better results
            if avg_confidence < 50:
                text_alt = pytesseract.image_to_string(image, config='--psm 3')
                if len(text_alt) > len(text):
                    text = text_alt
            
            return {
                'text': text.strip(),
                'confidence': avg_confidence,
                'method': 'tesseract'
            }
        
        except Exception as e:
            print(f"OCR error for {image_path}: {e}")
            return {'text': '', 'confidence': 0, 'method': 'failed'}
    
    def _analyze_image_content(self, image_path: str, ocr_text: str) -> Dict[str, Any]:
        """Analyze image content to determine type and characteristics"""
        analysis = {
            'type': 'unknown',
            'has_math': False,
            'has_code': False,
            'has_graph': False,
            'description': ''
        }
        
        ocr_lower = ocr_text.lower()
        
        # Check for mathematical content
        math_indicators = ['=', '+', '-', '×', '÷', '∑', '∫', '√', 'equation', 'formula']
        analysis['has_math'] = any(indicator in ocr_lower for indicator in math_indicators)
        
        # Check for code content
        code_indicators = ['def ', 'function', 'class', 'import', 'var ', 'const ', 'if (', 'for (']
        analysis['has_code'] = any(indicator in ocr_lower for indicator in code_indicators)
        
        # Check for graphs/charts
        graph_indicators = ['axis', 'chart', 'graph', 'plot', 'x-axis', 'y-axis', 'legend']
        analysis['has_graph'] = any(indicator in ocr_lower for indicator in graph_indicators)
        
        # Determine primary type
        if analysis['has_graph']:
            analysis['type'] = 'graph_chart'
        elif analysis['has_math']:
            analysis['type'] = 'mathematical'
        elif analysis['has_code']:
            analysis['type'] = 'code_snippet'
        elif len(ocr_text.strip()) > 20:
            analysis['type'] = 'text_image'
        else:
            analysis['type'] = 'diagram'
        
        # Create description
        if analysis['type'] == 'graph_chart':
            analysis['description'] = f"Graph or chart containing: {ocr_text[:100]}..."
        elif analysis['type'] == 'mathematical':
            analysis['description'] = f"Mathematical content: {ocr_text[:100]}..."
        elif analysis['type'] == 'code_snippet':
            analysis['description'] = f"Code snippet: {ocr_text[:100]}..."
        else:
            analysis['description'] = f"Image with text: {ocr_text[:100]}..."
        
        return analysis
    
    def _add_to_vector_db(self, processed_content: Dict[str, Any], filename: str) -> str:
        """Add processed content to vector database"""
        doc_id = str(uuid.uuid4())
        
        # Add text chunks
        if processed_content["text_chunks"]:
            self._add_text_chunks_to_db(processed_content["text_chunks"], doc_id, filename)
        
        # Add images
        if processed_content["images"]:
            self._add_images_to_db(processed_content["images"], doc_id, filename)
        
        # Add code blocks
        if processed_content["code_blocks"]:
            self._add_code_blocks_to_db(processed_content["code_blocks"], doc_id, filename)
        
        # Add formulas
        if processed_content["formulas"]:
            self._add_formulas_to_db(processed_content["formulas"], doc_id, filename)
        
        return doc_id
    
    def _add_text_chunks_to_db(self, text_chunks: List[Dict], doc_id: str, filename: str):
        """Add text chunks to vector database"""
        texts = [chunk["text"] for chunk in text_chunks]
        embeddings = self.embedding_model.encode(texts).tolist()
        
        ids = [f"{doc_id}_{chunk['chunk_id']}" for chunk in text_chunks]
        metadatas = []
        
        for chunk in text_chunks:
            metadata = {
                "document_id": doc_id,
                "filename": filename,
                "chunk_type": "text",
                "page_number": chunk.get("page_number", chunk.get("slide_number", 0)),
                "word_count": chunk.get("word_count", 0),
                "section_type": chunk.get("section_type", "general")
            }
            metadatas.append(metadata)
        
        self.text_collection.add(
            embeddings=embeddings,
            documents=texts,
            metadatas=metadatas,
            ids=ids
        )
    
    def _add_images_to_db(self, images: List[Dict], doc_id: str, filename: str):
        """Add images to vector database"""
        # Filter images with meaningful OCR text
        valid_images = [img for img in images if len(img.get("ocr_text", "").strip()) > 5]
        
        if not valid_images:
            return
        
        texts = [img["ocr_text"] for img in valid_images]
        embeddings = self.embedding_model.encode(texts).tolist()
        
        ids = [f"{doc_id}_{img['image_id']}" for img in valid_images]
        metadatas = []
        documents = []
        
        for img in valid_images:
            metadata = {
                "document_id": doc_id,
                "filename": filename,
                "chunk_type": "image",
                "page_number": img.get("page_number", img.get("slide_number", 0)),
                "image_path": img["image_path"],
                "image_type": img.get("image_type", "unknown"),
                "contains_math": img.get("contains_math", False),
                "contains_code": img.get("contains_code", False),
                "ocr_confidence": img.get("ocr_confidence", 0)
            }
            metadatas.append(metadata)
            documents.append(img["ocr_text"])
        
        self.image_collection.add(
            embeddings=embeddings,
            documents=documents,
            metadatas=metadatas,
            ids=ids
        )
    
    def _add_code_blocks_to_db(self, code_blocks: List[Dict], doc_id: str, filename: str):
        """Add code blocks to vector database"""
        codes = [block["code"] for block in code_blocks]
        embeddings = self.embedding_model.encode(codes).tolist()
        
        ids = [f"{doc_id}_{block['code_id']}" for block in code_blocks]
        metadatas = []
        
        for block in code_blocks:
            metadata = {
                "document_id": doc_id,
                "filename": filename,
                "chunk_type": "code",
                "page_number": block.get("page_number", block.get("slide_number", 0)),
                "language": block.get("language", "unknown"),
                "confidence": block.get("confidence", 0.5)
            }
            metadatas.append(metadata)
        
        self.code_collection.add(
            embeddings=embeddings,
            documents=codes,
            metadatas=metadatas,
            ids=ids
        )
    
    def _add_formulas_to_db(self, formulas: List[Dict], doc_id: str, filename: str):
        """Add formulas to vector database"""
        formula_texts = [formula["formula"] for formula in formulas]
        embeddings = self.embedding_model.encode(formula_texts).tolist()
        
        ids = [f"{doc_id}_{formula['formula_id']}" for formula in formulas]
        metadatas = []
        
        for formula in formulas:
            metadata = {
                "document_id": doc_id,
                "filename": filename,
                "chunk_type": "formula",
                "page_number": formula.get("page_number", formula.get("slide_number", 0)),
                "formula_type": formula.get("formula_type", "unknown"),
                "confidence": formula.get("confidence", 0.5)
            }
            metadatas.append(metadata)
        
        self.formula_collection.add(
            embeddings=embeddings,
            documents=formula_texts,
            metadatas=metadatas,
            ids=ids
        )
    
    def intelligent_search(self, query: str, top_k: int = 10) -> List[Dict[str, Any]]:
        """Intelligent search with content type weighting"""
        query_embedding = self.embedding_model.encode([query]).tolist()[0]
        
        all_results = []
        
        # Search each collection with different weights
        collections_config = [
            (self.text_collection, "text", 1.0),
            (self.image_collection, "image", 0.8),
            (self.code_collection, "code", 0.9),
            (self.formula_collection, "formula", 0.7)
        ]
        
        for collection, content_type, weight in collections_config:
            try:
                results = collection.query(
                    query_embeddings=[query_embedding],
                    n_results=min(top_k, 5)
                )
                
                for i in range(len(results['ids'][0])):
                    result = {
                        "id": results['ids'][0][i],
                        "content": results['documents'][0][i],
                        "metadata": results['metadatas'][0][i],
                        "distance": results['distances'][0][i],
                        "content_type": content_type,
                        "weighted_score": (1 - results['distances'][0][i]) * weight,
                        "relevance_score": self._calculate_relevance(query, results['documents'][0][i])
                    }
                    all_results.append(result)
            
            except Exception as e:
                print(f"Error searching {content_type} collection: {e}")
        
        # Sort by weighted score and relevance
        all_results.sort(key=lambda x: (x['weighted_score'] + x['relevance_score']) / 2, reverse=True)
        return all_results[:top_k]
    
    def _calculate_relevance(self, query: str, content: str) -> float:
        """Calculate semantic relevance between query and content"""
        query_words = set(query.lower().split())
        content_words = set(content.lower().split())
        
        # Jaccard similarity
        intersection = query_words.intersection(content_words)
        union = query_words.union(content_words)
        
        return len(intersection) / len(union) if union else 0
    
    def generate_prerequisites_with_ai(self, query: str, target_level: str = "undergraduate") -> Dict[str, Any]:
        """Generate prerequisites using AI with actual document content"""
        
        # Search for relevant content
        search_results = self.intelligent_search(query, top_k=20)
        
        if not search_results:
            return {
                "error": "No relevant content found in uploaded documents",
                "suggestions": "Please upload relevant educational materials first"
            }
        
        # Organize content by type
        content_by_type = {
            "text": [r for r in search_results if r["content_type"] == "text"],
            "images": [r for r in search_results if r["content_type"] == "image"],
            "code": [r for r in search_results if r["content_type"] == "code"],
            "formulas": [r for r in search_results if r["content_type"] == "formula"]
        }
        
        # Generate AI-powered prerequisite analysis
        prerequisite_analysis = self._analyze_prerequisites_with_ai(query, content_by_type, target_level)
        
        # Create structured prerequisite document
        prerequisite_doc = self._create_prerequisite_structure(prerequisite_analysis, content_by_type)
        
        return {
            "analysis": prerequisite_analysis,
            "structured_content": prerequisite_doc,
            "source_documents": list(set([r["metadata"]["filename"] for r in search_results])),
            "content_stats": {
                "text_chunks": len(content_by_type["text"]),
                "images": len(content_by_type["images"]),
                "code_examples": len(content_by_type["code"]),
                "formulas": len(content_by_type["formulas"])
            }
        }
    
    def _analyze_prerequisites_with_ai(self, query: str, content_by_type: Dict, target_level: str) -> Dict[str, Any]:
        """Use AI to analyze and identify prerequisites"""
        
        # Prepare context from retrieved content
        context_text = ""
        
        # Add text content
        for item in content_by_type["text"][:5]:
            context_text += f"Text: {item['content'][:500]}\n\n"
        
        # Add image descriptions
        for item in content_by_type["images"][:3]:
            context_text += f"Image content: {item['content'][:300]}\n\n"
        
        # Add code examples
        for item in content_by_type["code"][:3]:
            context_text += f"Code example: {item['content'][:400]}\n\n"
        
        # Add formulas
        for item in content_by_type["formulas"][:5]:
            context_text += f"Formula: {item['content']}\n\n"
        
        # Create AI prompt
        prompt = f"""
        As an educational expert, analyze the following content and identify prerequisite knowledge needed for the topic: "{query}"
        
        Target Level: {target_level}
        
        Content from uploaded documents:
        {context_text}
        
        Please provide a detailed analysis including:
        1. Fundamental concepts students should know before studying this topic
        2. Mathematical prerequisites (if any)
        3. Programming concepts (if applicable)
        4. Key terminology and definitions
        5. Suggested learning sequence
        
        Format your response as a structured analysis with clear sections.
        """
        
        try:
            response = self.model.generate_content(prompt)
            
            ai_analysis = response.text
            
            return {
                "ai_analysis": ai_analysis,
                "prerequisite_topics": self._extract_prerequisite_topics(ai_analysis),
                "learning_sequence": self._extract_learning_sequence(ai_analysis),
                "difficulty_level": self._assess_difficulty_level(ai_analysis, target_level)
            }
        
        except Exception as e:
            print(f"AI analysis error: {e}")
            # Fallback to rule-based analysis
            return self._fallback_prerequisite_analysis(query, content_by_type, target_level)
    
    def _extract_prerequisite_topics(self, ai_analysis: str) -> List[str]:
        """Extract prerequisite topics from AI analysis"""
        import re
        
        # Look for numbered lists, bullet points, or clear topic mentions
        topics = []
        
        # Pattern for numbered or bulleted lists
        list_pattern = r'(?:^\d+\.|\*|\-)\s*([^\n]+)'
        matches = re.findall(list_pattern, ai_analysis, re.MULTILINE)
        
        for match in matches:
            topic = match.strip()
            if len(topic) > 10 and not topic.startswith(('The', 'This', 'It', 'For')):
                topics.append(topic)
        
        # Also look for concepts mentioned after "prerequisite" or "required"
        prereq_pattern = r'(?:prerequisite|required|need|should know|must understand)[^.]*?([A-Z][^.]+)'
        prereq_matches = re.findall(prereq_pattern, ai_analysis, re.IGNORECASE)
        
        for match in prereq_matches:
            if len(match.strip()) > 10:
                topics.append(match.strip())
        
        return list(set(topics))[:10]  # Return unique topics, max 10
    
    def _extract_learning_sequence(self, ai_analysis: str) -> List[str]:
        """Extract suggested learning sequence from AI analysis"""
        import re
        
        sequence = []
        
        # Look for sequence indicators
        sequence_patterns = [
            r'(?:first|start|begin)[^.]*?([A-Z][^.]+)',
            r'(?:then|next|after)[^.]*?([A-Z][^.]+)',
            r'(?:finally|last|end)[^.]*?([A-Z][^.]+)'
        ]
        
        for pattern in sequence_patterns:
            matches = re.findall(pattern, ai_analysis, re.IGNORECASE)
            for match in matches:
                if len(match.strip()) > 10:
                    sequence.append(match.strip())
        
        return sequence
    
    def _assess_difficulty_level(self, ai_analysis: str, target_level: str) -> Dict[str, Any]:
        """Assess difficulty level based on AI analysis"""
        difficulty_indicators = {
            "beginner": ["basic", "introduction", "fundamental", "simple"],
            "intermediate": ["intermediate", "moderate", "requires understanding"],
            "advanced": ["advanced", "complex", "sophisticated", "deep understanding"],
            "expert": ["expert", "research", "cutting-edge", "theoretical"]
        }
        
        level_scores = {}
        for level, indicators in difficulty_indicators.items():
            score = sum(1 for indicator in indicators if indicator in ai_analysis.lower())
            level_scores[level] = score
        
        assessed_level = max(level_scores, key=level_scores.get)
        
        return {
            "assessed_level": assessed_level,
            "target_level": target_level,
            "level_match": assessed_level == target_level,
            "complexity_score": level_scores[assessed_level]
        }
    
    def _fallback_prerequisite_analysis(self, query: str, content_by_type: Dict, target_level: str) -> Dict[str, Any]:
        """Fallback analysis when AI is not available"""
        
        # Rule-based prerequisite identification
        prerequisite_topics = []
        
        # Analyze text content for concepts
        for item in content_by_type["text"]:
            content = item["content"].lower()
            
            # Look for definition patterns
            if "definition" in content or "define" in content:
                # Extract the concept being defined
                import re
                def_matches = re.findall(r'(?:definition of|define)\s+([^.]+)', content)
                prerequisite_topics.extend([match.strip() for match in def_matches])
        
        # Analyze mathematical content
        math_topics = []
        for item in content_by_type["formulas"]:
            formula = item["content"]
            if any(symbol in formula for symbol in ['∑', '∫', '∂', 'lim']):
                math_topics.append("Calculus")
            elif any(symbol in formula for symbol in ['σ', 'μ', 'P(', 'E[']):
                math_topics.append("Statistics and Probability")
            elif 'matrix' in formula.lower() or 'vector' in formula.lower():
                math_topics.append("Linear Algebra")
        
        prerequisite_topics.extend(math_topics)
        
        # Analyze code content
        programming_topics = []
        for item in content_by_type["code"]:
            code = item["content"]
            language = item["metadata"].get("language", "unknown")
            
            if language == "python":
                programming_topics.append("Python Programming Basics")
            elif language == "javascript":
                programming_topics.append("JavaScript Fundamentals")
            elif "class" in code.lower():
                programming_topics.append("Object-Oriented Programming")
            elif "function" in code.lower() or "def " in code:
                programming_topics.append("Functions and Methods")
        
        prerequisite_topics.extend(programming_topics)
        
        return {
            "ai_analysis": "Rule-based analysis (AI not available)",
            "prerequisite_topics": list(set(prerequisite_topics)),
            "learning_sequence": prerequisite_topics,
            "difficulty_level": {
                "assessed_level": target_level,
                "target_level": target_level,
                "level_match": True,
                "complexity_score": 1
            }
        }
    
    def _create_prerequisite_structure(self, analysis: Dict, content_by_type: Dict) -> Dict[str, Any]:
        """Create structured prerequisite document with actual content"""
        
        structured_doc = {
            "title": "Prerequisite Knowledge Document",
            "sections": {}
        }
        
        # Create sections based on content types and analysis
        
        # 1. Fundamental Concepts (from text)
        if content_by_type["text"]:
            fundamental_concepts = []
            for item in content_by_type["text"][:5]:
                concept = {
                    "title": self._extract_concept_title(item["content"]),
                    "description": item["content"][:300] + "...",
                    "source": item["metadata"]["filename"],
                    "page": item["metadata"]["page_number"]
                }
                fundamental_concepts.append(concept)
            
            structured_doc["sections"]["Fundamental Concepts"] = {
                "type": "text",
                "content": fundamental_concepts
            }
        
        # 2. Mathematical Prerequisites (from formulas)
        if content_by_type["formulas"]:
            math_prereqs = []
            for item in content_by_type["formulas"]:
                formula = {
                    "formula": item["content"],
                    "explanation": f"Mathematical concept extracted from {item['metadata']['filename']}",
                    "source": item["metadata"]["filename"],
                    "page": item["metadata"]["page_number"]
                }
                math_prereqs.append(formula)
            
            structured_doc["sections"]["Mathematical Prerequisites"] = {
                "type": "formulas",
                "content": math_prereqs
            }
        
        # 3. Visual Content (from images)
        if content_by_type["images"]:
            visual_content = []
            for item in content_by_type["images"]:
                image_info = {
                    "image_path": item["metadata"]["image_path"],
                    "description": item["content"],
                    "type": item["metadata"].get("image_type", "diagram"),
                    "source": item["metadata"]["filename"],
                    "page": item["metadata"]["page_number"]
                }
                visual_content.append(image_info)
            
            structured_doc["sections"]["Visual Learning Materials"] = {
                "type": "images",
                "content": visual_content
            }
        
        # 4. Code Examples (from code)
        if content_by_type["code"]:
            code_examples = []
            for item in content_by_type["code"]:
                code_info = {
                    "code": item["content"],
                    "language": item["metadata"].get("language", "unknown"),
                    "explanation": f"Code example from {item['metadata']['filename']}",
                    "source": item["metadata"]["filename"],
                    "page": item["metadata"]["page_number"]
                }
                code_examples.append(code_info)
            
            structured_doc["sections"]["Programming Examples"] = {
                "type": "code",
                "content": code_examples
            }
        
        return structured_doc
    
    def _extract_concept_title(self, text: str) -> str:
        """Extract a title/concept name from text content"""
        # Look for the first sentence or a clear concept name
        sentences = text.split('.')
        first_sentence = sentences[0].strip()
        
        # If it's a definition, extract the concept
        if "definition" in first_sentence.lower():
            import re
            match = re.search(r'definition of ([^:]+)', first_sentence, re.IGNORECASE)
            if match:
                return match.group(1).strip()
        
        # Otherwise, return first few words
        words = first_sentence.split()[:5]
        return ' '.join(words) + ("..." if len(words) == 5 else "")
    
    def generate_chat_response(self, query: str, conversation_history: List[Dict] = None) -> str:
        """Generate conversational response based on uploaded documents"""
        
        # Search for relevant content
        search_results = self.intelligent_search(query, top_k=5)
        
        if not search_results:
            return "I don't have information about that topic in the uploaded documents. Please upload relevant materials first."
        
        # Prepare context
        context = ""
        images_info = []
        
        for result in search_results:
            if result["content_type"] == "text":
                context += f"From {result['metadata']['filename']}: {result['content'][:300]}...\n\n"
            elif result["content_type"] == "image":
                images_info.append({
                    "path": result["metadata"]["image_path"],
                    "description": result["content"],
                    "source": result["metadata"]["filename"]
                })
                context += f"Image content from {result['metadata']['filename']}: {result['content'][:200]}...\n\n"
            elif result["content_type"] == "code":
                context += f"Code example from {result['metadata']['filename']}:\n{result['content'][:400]}\n\n"
            elif result["content_type"] == "formula":
                context += f"Formula from {result['metadata']['filename']}: {result['content']}\n\n"
        
        # Generate AI response
        try:
            prompt = f"""
            Based on the following content from uploaded educational documents, answer the question: "{query}"
            
            Context from documents:
            {context}
            
            Please provide a comprehensive answer that:
            1. Directly addresses the question
            2. References the source materials
            3. Includes relevant examples from the documents
            4. Mentions any images, graphs, or diagrams if relevant
            
            Keep the response educational and well-structured.
            """
            
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are an AI tutor helping students learn from their uploaded educational materials."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=800,
                temperature=0.4
            )
            
            ai_response = response.choices[0].message.content
            
            # Add image information if available
            if images_info:
                ai_response += "\n\n**Related Visual Materials:**\n"
                for img in images_info[:3]:  # Limit to 3 images
                    ai_response += f"- {img['description'][:100]}... (from {img['source']})\n"
            
            return ai_response
        
        except Exception as e:
            print(f"AI response error: {e}")
            # Fallback response
            return f"Based on your uploaded documents, here's what I found:\n\n{context[:500]}..."
    
    def get_document_stats(self) -> Dict[str, Any]:
        """Get comprehensive statistics about processed documents"""
        stats = {
            "total_documents": 0,
            "content_counts": {
                "text_chunks": 0,
                "images": 0,
                "code_blocks": 0,
                "formulas": 0
            },
            "by_document": {},
            "processing_info": {
                "embedding_model": str(self.embedding_model),
                "vector_db_path": self.vector_db_dir,
                "images_extracted": len(os.listdir(self.images_dir)) if os.path.exists(self.images_dir) else 0
            }
        }
        
        # Get document-specific stats
        try:
            # Get all documents from text collection
            text_results = self.text_collection.get()
            doc_files = {}
            
            for metadata in text_results.get('metadatas', []):
                filename = metadata.get('filename', 'unknown')
                doc_id = metadata.get('document_id', 'unknown')
                
                if filename not in doc_files:
                    doc_files[filename] = {
                        "document_id": doc_id,
                        "text_chunks": 0,
                        "images": 0,
                        "code_blocks": 0,
                        "formulas": 0
                    }
                
                doc_files[filename]["text_chunks"] += 1
            
            # Count other content types
            collections = [
                (self.image_collection, "images"),
                (self.code_collection, "code_blocks"),
                (self.formula_collection, "formulas")
            ]
            
            for collection, content_type in collections:
                try:
                    results = collection.get()
                    for metadata in results.get('metadatas', []):
                        filename = metadata.get('filename', 'unknown')
                        if filename in doc_files:
                            doc_files[filename][content_type] += 1
                        stats["content_counts"][content_type] += 1
                except Exception as e:
                    print(f"Error counting {content_type}: {e}")
            
            stats["by_document"] = doc_files
            stats["total_documents"] = len(doc_files)
            stats["content_counts"]["text_chunks"] = sum(doc["text_chunks"] for doc in doc_files.values())
            
        except Exception as e:
            print(f"Error getting document stats: {e}")
        
        return stats