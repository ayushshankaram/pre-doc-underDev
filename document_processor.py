import os
import re
import json
import fitz  # PyMuPDF
from pptx import Presentation
from PIL import Image
import pytesseract
from typing import Dict, List, Any
import hashlib

class DocumentProcessor:
    def __init__(self, upload_dir: str, processed_dir: str):
        self.upload_dir = upload_dir
        self.processed_dir = processed_dir
        
    def process_document(self, file_path: str) -> Dict[str, Any]:
        """Process a document and extract all content types"""
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext == '.pdf':
            return self._process_pdf(file_path)
        elif file_ext in ['.ppt', '.pptx']:
            return self._process_ppt(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")
    
    def _process_pdf(self, file_path: str) -> Dict[str, Any]:
        """Process PDF document"""
        doc = fitz.open(file_path)
        
        result = {
            "filename": os.path.basename(file_path),
            "file_type": "pdf",
            "text_chunks": [],
            "images": [],
            "formulas": [],
            "code_blocks": [],
            "metadata": {
                "total_pages": len(doc),
                "file_size": os.path.getsize(file_path)
            }
        }
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # Extract text
            text = page.get_text()
            if text.strip():
                # Split into chunks by paragraphs
                chunks = self._split_into_chunks(text, page_num + 1)
                result["text_chunks"].extend(chunks)
                
                # Extract code blocks
                code_blocks = self._extract_code_blocks(text, page_num + 1)
                result["code_blocks"].extend(code_blocks)
                
                # Extract formulas (basic pattern matching)
                formulas = self._extract_formulas(text, page_num + 1)
                result["formulas"].extend(formulas)
            
            # Extract images
            images = self._extract_images_from_pdf_page(page, page_num + 1)
            result["images"].extend(images)
        
        doc.close()
        
        # Save processed content
        self._save_processed_content(result)
        
        return result
    
    def _process_ppt(self, file_path: str) -> Dict[str, Any]:
        """Process PowerPoint document"""
        prs = Presentation(file_path)
        
        result = {
            "filename": os.path.basename(file_path),
            "file_type": "pptx",
            "text_chunks": [],
            "images": [],
            "formulas": [],
            "code_blocks": [],
            "metadata": {
                "total_slides": len(prs.slides),
                "file_size": os.path.getsize(file_path)
            }
        }
        
        for slide_num, slide in enumerate(prs.slides, 1):
            # Extract text from slide
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            
            if slide_text.strip():
                # Create slide-level chunks
                chunks = self._split_into_chunks(slide_text, slide_num, is_slide=True)
                result["text_chunks"].extend(chunks)
                
                # Extract code blocks
                code_blocks = self._extract_code_blocks(slide_text, slide_num, is_slide=True)
                result["code_blocks"].extend(code_blocks)
                
                # Extract formulas
                formulas = self._extract_formulas(slide_text, slide_num, is_slide=True)
                result["formulas"].extend(formulas)
            
            # Extract images from slide
            images = self._extract_images_from_slide(slide, slide_num)
            result["images"].extend(images)
        
        # Save processed content
        self._save_processed_content(result)
        
        return result
    
    def _split_into_chunks(self, text: str, page_num: int, is_slide: bool = False) -> List[Dict[str, Any]]:
        """Split text into meaningful chunks"""
        chunks = []
        
        # Split by paragraphs first
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        
        for i, paragraph in enumerate(paragraphs):
            if len(paragraph) < 50:  # Skip very short paragraphs
                continue
                
            chunk = {
                "text": paragraph,
                "chunk_id": f"{'slide' if is_slide else 'page'}_{page_num}_chunk_{i}",
                "page_number" if not is_slide else "slide_number": page_num,
                "chunk_type": "text",
                "word_count": len(paragraph.split()),
                "char_count": len(paragraph)
            }
            chunks.append(chunk)
        
        return chunks
    
    def _extract_code_blocks(self, text: str, page_num: int, is_slide: bool = False) -> List[Dict[str, Any]]:
        """Extract code blocks from text"""
        code_blocks = []
        
        # Pattern for code blocks (indented text or code-like patterns)
        code_patterns = [
            r'```[\s\S]*?```',  # Markdown code blocks
            r'def\s+\w+\([^)]*\):[\s\S]*?(?=\n\S|\n\n|$)',  # Python functions
            r'class\s+\w+[\s\S]*?(?=\n\S|\n\n|$)',  # Python classes
            r'for\s+\w+\s+in\s+[\s\S]*?(?=\n\S|\n\n|$)',  # For loops
            r'if\s+[\s\S]*?(?=\n\S|\n\n|$)',  # If statements
        ]
        
        for i, pattern in enumerate(code_patterns):
            matches = re.finditer(pattern, text, re.MULTILINE)
            for match in matches:
                code_block = {
                    "code": match.group().strip(),
                    "code_id": f"{'slide' if is_slide else 'page'}_{page_num}_code_{i}",
                    "page_number" if not is_slide else "slide_number": page_num,
                    "language": self._detect_language(match.group()),
                    "chunk_type": "code"
                }
                code_blocks.append(code_block)
        
        return code_blocks
    
    def _extract_formulas(self, text: str, page_num: int, is_slide: bool = False) -> List[Dict[str, Any]]:
        """Extract mathematical formulas from text"""
        formulas = []
        
        # Simple patterns for mathematical expressions
        formula_patterns = [
            r'\$[^$]+\$',  # LaTeX inline math
            r'\$\$[^$]+\$\$',  # LaTeX display math
            r'[∑∏∫∆∇∞≈≠≤≥±×÷√∠∴∵∪∩∈∀∃]',  # Mathematical symbols
            r'\b\w+\s*=\s*[^,\n]+',  # Equations
        ]
        
        for i, pattern in enumerate(formula_patterns):
            matches = re.finditer(pattern, text)
            for match in matches:
                formula = {
                    "formula": match.group().strip(),
                    "formula_id": f"{'slide' if is_slide else 'page'}_{page_num}_formula_{i}",
                    "page_number" if not is_slide else "slide_number": page_num,
                    "chunk_type": "formula"
                }
                formulas.append(formula)
        
        return formulas
    
    def _extract_images_from_pdf_page(self, page, page_num: int) -> List[Dict[str, Any]]:
        """Extract images from PDF page"""
        images = []
        image_list = page.get_images()
        
        for img_index, img in enumerate(image_list):
            try:
                xref = img[0]
                pix = fitz.Pixmap(page.parent, xref)
                
                if pix.n - pix.alpha < 4:  # Skip if not RGB
                    img_data = pix.tobytes("png")
                    
                    # Save image
                    img_filename = f"page_{page_num}_img_{img_index}.png"
                    img_path = os.path.join(self.processed_dir, img_filename)
                    
                    with open(img_path, "wb") as img_file:
                        img_file.write(img_data)
                    
                    # OCR the image
                    ocr_text = self._ocr_image(img_path)
                    
                    image_info = {
                        "image_path": img_path,
                        "image_id": f"page_{page_num}_img_{img_index}",
                        "page_number": page_num,
                        "ocr_text": ocr_text,
                        "chunk_type": "image",
                        "width": pix.width,
                        "height": pix.height
                    }
                    images.append(image_info)
                
                pix = None
            except Exception as e:
                print(f"Error extracting image {img_index} from page {page_num}: {e}")
        
        return images
    
    def _extract_images_from_slide(self, slide, slide_num: int) -> List[Dict[str, Any]]:
        """Extract images from PowerPoint slide"""
        images = []
        
        for shape_index, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # Picture shape type
                try:
                    image = shape.image
                    img_filename = f"slide_{slide_num}_img_{shape_index}.png"
                    img_path = os.path.join(self.processed_dir, img_filename)
                    
                    with open(img_path, "wb") as img_file:
                        img_file.write(image.blob)
                    
                    # OCR the image
                    ocr_text = self._ocr_image(img_path)
                    
                    image_info = {
                        "image_path": img_path,
                        "image_id": f"slide_{slide_num}_img_{shape_index}",
                        "slide_number": slide_num,
                        "ocr_text": ocr_text,
                        "chunk_type": "image"
                    }
                    images.append(image_info)
                
                except Exception as e:
                    print(f"Error extracting image {shape_index} from slide {slide_num}: {e}")
        
        return images
    
    def _ocr_image(self, image_path: str) -> str:
        """Perform OCR on image"""
        try:
            image = Image.open(image_path)
            text = pytesseract.image_to_string(image)
            return text.strip()
        except Exception as e:
            print(f"Error performing OCR on {image_path}: {e}")
            return ""
    
    def _detect_language(self, code: str) -> str:
        """Detect programming language of code block"""
        # Simple heuristics for language detection
        if 'def ' in code and 'import ' in code:
            return 'python'
        elif 'function' in code and '{' in code:
            return 'javascript'
        elif 'public class' in code:
            return 'java'
        elif '#include' in code:
            return 'cpp'
        else:
            return 'unknown'
    
    def _save_processed_content(self, content: Dict[str, Any]):
        """Save processed content to file"""
        filename = content["filename"]
        processed_filename = f"processed_{filename}.json"
        processed_path = os.path.join(self.processed_dir, processed_filename)
        
        with open(processed_path, 'w', encoding='utf-8') as f:
            json.dump(content, f, indent=2, ensure_ascii=False)
    
    def get_document_hash(self, file_path: str) -> str:
        """Generate hash for document to detect duplicates"""
        with open(file_path, 'rb') as f:
            content = f.read()
            return hashlib.md5(content).hexdigest()